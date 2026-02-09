"""
Navigation helpers for PowerPoint presentations.

Provides utilities for:
- Creating Table of Contents slides with hyperlinks (multi-column, multi-slide)
- Adding PowerPoint sections for navigation
- Creating internal hyperlinks between slides
"""

from typing import List, Dict, Optional, Tuple, TYPE_CHECKING
import uuid
import math

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree

if TYPE_CHECKING:
    from pptx import Presentation
    from pptx.slide import Slide

from utilities.pptx_helper.layouts import CONTENT, BLANK

# TOC layout constants
TOC_MAX_ENTRIES_PER_COLUMN = 15  # Maximum entries per column before overflow
TOC_MAX_COLUMNS = 3  # Maximum columns per slide
TOC_FONT_SIZE_LARGE = Pt(16)  # Font size for fewer entries
TOC_FONT_SIZE_MEDIUM = Pt(14)  # Font size for medium number of entries
TOC_FONT_SIZE_SMALL = Pt(12)  # Font size for many entries
TOC_LINE_SPACING = 1.1  # Line spacing multiplier


def add_internal_hyperlink(run, source_slide: "Slide", target_slide: "Slide") -> str:
    """
    Add an internal hyperlink from a text run to a target slide.
    
    Args:
        run: The text run to make into a hyperlink
        source_slide: The slide containing the hyperlink
        target_slide: The slide to link to
        
    Returns:
        The relationship ID created
    """
    # Create relationship from source slide to target slide
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    
    # Get or create rPr element
    rPr = run._r.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.SubElement(run._r, qn('a:rPr'))
        run._r.insert(0, rPr)
    
    # Remove any existing hlinkClick
    for existing in rPr.findall(qn('a:hlinkClick')):
        rPr.remove(existing)
    
    # Create hlinkClick element with action for internal link
    hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')
    
    return rId


def add_section(prs: "Presentation", name: str, slide_ids: List[int]) -> None:
    """
    Add a PowerPoint section containing the specified slides.
    
    Args:
        prs: Presentation object
        name: Section name
        slide_ids: List of slide IDs to include in the section
    """
    prs_elm = prs.part._element
    
    # Find or create sectionLst
    sectionLst = prs_elm.find(qn('p:sectionLst'))
    if sectionLst is None:
        sldIdLst = prs_elm.find(qn('p:sldIdLst'))
        if sldIdLst is not None:
            sectionLst = etree.Element(qn('p:sectionLst'))
            sldIdLst.addnext(sectionLst)
    
    if sectionLst is None:
        return
    
    # Create section element
    section = etree.SubElement(sectionLst, qn('p:section'))
    section.set('name', name)
    section.set('id', '{' + str(uuid.uuid4()).upper() + '}')
    
    # Add slide IDs to section
    sec_sldIdLst = etree.SubElement(section, qn('p:sldIdLst'))
    for slide_id in slide_ids:
        sldId = etree.SubElement(sec_sldIdLst, qn('p:sldId'))
        sldId.set('id', str(slide_id))


def create_sections_from_headers(prs: "Presentation", section_slides: Dict[str, List["Slide"]]) -> None:
    """
    Create PowerPoint sections from a dictionary of section names to slides.
    
    Args:
        prs: Presentation object
        section_slides: Dict mapping section names to lists of slides
    """
    for section_name, slides in section_slides.items():
        slide_ids = [slide.slide_id for slide in slides]
        add_section(prs, section_name, slide_ids)


def _calculate_toc_layout(num_entries: int) -> Tuple[int, int, int]:
    """
    Calculate optimal TOC layout based on number of entries.
    
    Returns:
        Tuple of (num_columns, entries_per_column, num_slides)
    """
    # Determine number of columns needed (prefer 2 columns for readability)
    if num_entries <= TOC_MAX_ENTRIES_PER_COLUMN:
        # Single column fits
        return (1, num_entries, 1)
    elif num_entries <= TOC_MAX_ENTRIES_PER_COLUMN * 2:
        # Two columns fit on one slide
        entries_per_col = math.ceil(num_entries / 2)
        return (2, entries_per_col, 1)
    elif num_entries <= TOC_MAX_ENTRIES_PER_COLUMN * 3:
        # Three columns fit on one slide
        entries_per_col = math.ceil(num_entries / 3)
        return (3, entries_per_col, 1)
    else:
        # Need multiple slides with 3 columns each
        entries_per_slide = TOC_MAX_ENTRIES_PER_COLUMN * 3
        num_slides = math.ceil(num_entries / entries_per_slide)
        return (3, TOC_MAX_ENTRIES_PER_COLUMN, num_slides)


def _get_toc_font_size(num_entries: int, num_columns: int) -> Pt:
    """Get appropriate font size based on content density."""
    entries_per_col = math.ceil(num_entries / num_columns) if num_columns > 0 else num_entries
    
    if entries_per_col <= 10:
        return TOC_FONT_SIZE_LARGE
    elif entries_per_col <= 15:
        return TOC_FONT_SIZE_MEDIUM
    else:
        return TOC_FONT_SIZE_SMALL


def _create_toc_column(
    slide: "Slide",
    entries: List[Tuple[str, "Slide", int]],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: Pt
) -> None:
    """Create a single column of TOC entries."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    
    for i, (entry_text, target_slide, indent_level) in enumerate(entries):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.level = indent_level
        p.space_after = Pt(2)  # Tighter spacing
        
        run = p.add_run()
        
        # Truncate long entries to fit column width
        max_chars = int(width * 8)  # Approximate chars that fit
        if len(entry_text) > max_chars:
            entry_text = entry_text[:max_chars-3] + "..."
        
        run.text = entry_text
        run.font.size = font_size
        run.font.bold = (indent_level == 0)
        
        # Style as hyperlink
        run.font.color.rgb = RGBColor(0, 51, 153)  # Dark blue
        run.font.underline = True
        
        # Add internal hyperlink
        add_internal_hyperlink(run, slide, target_slide)


def add_toc_slide(
    prs: "Presentation",
    title: str,
    toc_entries: List[Tuple[str, "Slide", int]],
    subtitle: str = ""
) -> "Slide":
    """
    Add Table of Contents slide(s) with hyperlinks to sections.
    
    Automatically uses multi-column layout and multiple slides if needed
    to prevent content overflow.
    
    Args:
        prs: Presentation object
        title: TOC slide title
        toc_entries: List of (entry_text, target_slide, indent_level) tuples
        subtitle: Optional subtitle
        
    Returns:
        The first created TOC slide (for positioning purposes)
    """
    slides = add_toc_slides(prs, title, toc_entries, subtitle)
    return slides[0] if slides else None


def add_toc_slides(
    prs: "Presentation",
    title: str,
    toc_entries: List[Tuple[str, "Slide", int]],
    subtitle: str = ""
) -> List["Slide"]:
    """
    Add Table of Contents slide(s) with hyperlinks to sections.
    
    Returns all created TOC slides (useful when multiple slides are needed).
    
    Args:
        prs: Presentation object
        title: TOC slide title
        toc_entries: List of (entry_text, target_slide, indent_level) tuples
        subtitle: Optional subtitle
        
    Returns:
        List of all created TOC slides
    """
    if not toc_entries:
        layout = prs.slide_layouts[CONTENT]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        return [slide]
    
    num_entries = len(toc_entries)
    num_columns, entries_per_column, num_slides = _calculate_toc_layout(num_entries)
    font_size = _get_toc_font_size(num_entries, num_columns)
    
    # Content area dimensions (in inches)
    content_left = 0.5
    content_top = 1.3
    content_width = 12.3
    content_height = 5.5
    
    # Calculate column dimensions (in inches)
    column_gap = 0.3
    total_gap = column_gap * (num_columns - 1) if num_columns > 1 else 0
    column_width = (content_width - total_gap) / num_columns
    
    slides_created = []
    entry_index = 0
    
    for slide_num in range(num_slides):
        layout = prs.slide_layouts[CONTENT]
        slide = prs.slides.add_slide(layout)
        slides_created.append(slide)
        
        # Set title (with page number if multiple slides)
        if slide.shapes.title:
            if num_slides > 1:
                slide.shapes.title.text = f"{title} ({slide_num + 1}/{num_slides})"
            else:
                slide.shapes.title.text = title
        
        # Create columns for this slide
        for col in range(num_columns):
            if entry_index >= num_entries:
                break
            
            # Calculate entries for this column
            remaining = num_entries - entry_index
            entries_this_col = min(entries_per_column, remaining)
            
            # For last slide, distribute remaining entries evenly
            if slide_num == num_slides - 1 and col == 0:
                remaining_entries = num_entries - entry_index
                remaining_cols = num_columns
                entries_this_col = math.ceil(remaining_entries / remaining_cols)
            
            column_entries = toc_entries[entry_index:entry_index + entries_this_col]
            
            # Calculate column position (in inches)
            col_left = content_left + col * (column_width + column_gap)
            
            _create_toc_column(
                slide,
                column_entries,
                col_left,
                content_top,
                column_width,
                content_height,
                font_size
            )
            
            entry_index += entries_this_col
    
    return slides_created


def build_toc_from_yaml_sections(
    prs: "Presentation",
    sections: List[Dict],
    slide_map: Dict[int, "Slide"]
) -> List[Tuple[str, "Slide", int]]:
    """
    Build TOC entries from YAML sections data.
    
    Args:
        prs: Presentation object
        sections: List of section dictionaries from YAML
        slide_map: Dict mapping section index to slide object
        
    Returns:
        List of (entry_text, target_slide, indent_level) tuples
    """
    toc_entries = []
    current_section = None
    
    for i, section in enumerate(sections):
        section_type = section.get('type', '')
        
        if section_type == 'section_header':
            title = section.get('title', '')
            if i in slide_map:
                toc_entries.append((title, slide_map[i], 0))
                current_section = title
    
    return toc_entries


def insert_slide_at_position(prs: "Presentation", slide: "Slide", position: int) -> None:
    """
    Move a slide to a specific position in the presentation.
    
    Note: This manipulates the XML directly as python-pptx doesn't support
    slide reordering natively.
    
    Args:
        prs: Presentation object
        slide: Slide to move
        position: Target position (0-indexed)
    """
    # Get the slide's rId
    slide_rId = None
    for rel in prs.part.rels.values():
        if rel.target_part == slide.part:
            slide_rId = rel.rId
            break
    
    if slide_rId is None:
        return
    
    # Find the sldId element for this slide
    prs_elm = prs.part._element
    sldIdLst = prs_elm.find(qn('p:sldIdLst'))
    
    if sldIdLst is None:
        return
    
    # Find and remove the sldId
    target_sldId = None
    for sldId in sldIdLst.findall(qn('p:sldId')):
        if sldId.get(qn('r:id')) == slide_rId:
            target_sldId = sldId
            sldIdLst.remove(sldId)
            break
    
    if target_sldId is None:
        return
    
    # Insert at new position
    sldIds = list(sldIdLst.findall(qn('p:sldId')))
    if position >= len(sldIds):
        sldIdLst.append(target_sldId)
    else:
        sldIds[position].addprevious(target_sldId)


def insert_slides_at_position(prs: "Presentation", slides: List["Slide"], position: int) -> None:
    """
    Move multiple slides to a specific position in the presentation.
    
    Args:
        prs: Presentation object
        slides: List of slides to move (will be inserted in order)
        position: Target position for first slide (0-indexed)
    """
    # Insert in reverse order so they end up in correct order
    for i, slide in enumerate(reversed(slides)):
        insert_slide_at_position(prs, slide, position)
