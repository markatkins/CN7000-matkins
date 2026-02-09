"""
Slide layout constants and helpers for Cornelis PowerPoint template.

Usage:
    from utilities.pptx_helper.layouts import LAYOUTS, get_layout, CONTENT, TWO_COLUMN
    
    # Get layout by name
    layout = get_layout(prs, "content")
    
    # Or use constant
    layout = prs.slide_layouts[CONTENT]

Design:
    Layout indices are based on the Standard PPT Template_Light.potx template.
    Only the most commonly used layouts are exposed as constants.
    See templates/template-spec.md for complete layout documentation.
"""

from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from pptx import Presentation
    from pptx.slide import SlideLayout

# Layout index constants for key layouts
COVER = 0                    # Cover with motion - title slides with presenter info
COVER_PLAIN = 1              # Cover Plain presenter on right
COVER_B = 3                  # Cover B - alternative cover
BLANK = 4                    # Blank Layout - custom content
BLANK_WITH_LOGO = 5          # Blank with logo
TITLE_ONLY = 6               # 1_Title Only
CONTENT = 7                  # title and 1 content - standard content slide
CONTENT_WITH_SUBTITLE = 8    # Title subtitle and content
TWO_COLUMN = 12              # Title and Content - 2 Column - comparisons
IMAGE_LEFT = 16              # 1/3 Image Left
IMAGE_RIGHT = 20             # Image Right 2:3
HALF_IMAGE_RIGHT = 21        # Half Image Right
HALF_IMAGE_LEFT = 22         # Half Image Left
CONTENT_AND_PICTURE = 24     # 1/2 Content and Picture Insert
TABLE = 29                   # Title details and Table
KEYWORD = 30                 # Keyword Slide - emphasis
IMAGE_KEYWORD = 31           # Image Keyword Slide - full image
SECTION_HEADER = 36          # Section Header - section dividers

# Layout name to index mapping
LAYOUTS = {
    "cover": COVER,
    "cover_plain": COVER_PLAIN,
    "cover_b": COVER_B,
    "blank": BLANK,
    "blank_with_logo": BLANK_WITH_LOGO,
    "title_only": TITLE_ONLY,
    "content": CONTENT,
    "content_with_subtitle": CONTENT_WITH_SUBTITLE,
    "two_column": TWO_COLUMN,
    "image_left": IMAGE_LEFT,
    "image_right": IMAGE_RIGHT,
    "half_image_right": HALF_IMAGE_RIGHT,
    "half_image_left": HALF_IMAGE_LEFT,
    "content_and_picture": CONTENT_AND_PICTURE,
    "table": TABLE,
    "keyword": KEYWORD,
    "image_keyword": IMAGE_KEYWORD,
    "section_header": SECTION_HEADER,
}

# Placeholder indices for common layouts
# Format: {layout_name: {placeholder_name: idx}}
PLACEHOLDERS = {
    "cover": {
        "title": 0,
        "subtitle": 1,
        "presenter_name": 11,
        "presenter_info": 12,
        "presenter_name_2": 13,
    },
    "content": {
        "title": 0,
        "body": 1,
    },
    "content_with_subtitle": {
        "title": 0,
        "subtitle": 13,
        "body": 1,
    },
    "two_column": {
        "title": 0,
        "left": 1,
        "right": 11,
    },
    "content_and_picture": {
        "title": 0,
        "body": 1,
        "picture": 13,
    },
    "table": {
        "title": 0,
        "description": 10,
        "table": 11,
    },
    "section_header": {
        "title": 0,
        "subtitle": 1,
    },
}


def get_layout(prs: "Presentation", name: str) -> "SlideLayout":
    """
    Get a slide layout by name.
    
    Args:
        prs: The Presentation object
        name: Layout name (e.g., "content", "two_column", "table")
        
    Returns:
        SlideLayout object
        
    Raises:
        KeyError: If layout name is not recognized
        IndexError: If layout index is out of range
        
    Example:
        >>> prs = create_presentation()
        >>> layout = get_layout(prs, "content")
        >>> slide = prs.slides.add_slide(layout)
    """
    if name not in LAYOUTS:
        available = ", ".join(sorted(LAYOUTS.keys()))
        raise KeyError(f"Unknown layout '{name}'. Available: {available}")
    
    idx = LAYOUTS[name]
    return prs.slide_layouts[idx]


def get_placeholder_idx(layout_name: str, placeholder_name: str) -> int:
    """
    Get the placeholder index for a specific layout and placeholder.
    
    Args:
        layout_name: Name of the layout (e.g., "content", "two_column")
        placeholder_name: Name of the placeholder (e.g., "title", "body", "left")
        
    Returns:
        Placeholder index
        
    Raises:
        KeyError: If layout or placeholder name is not recognized
        
    Example:
        >>> idx = get_placeholder_idx("two_column", "right")
        >>> slide.placeholders[idx].text = "Right column content"
    """
    if layout_name not in PLACEHOLDERS:
        available = ", ".join(sorted(PLACEHOLDERS.keys()))
        raise KeyError(f"Unknown layout '{layout_name}'. Available: {available}")
    
    layout_placeholders = PLACEHOLDERS[layout_name]
    if placeholder_name not in layout_placeholders:
        available = ", ".join(sorted(layout_placeholders.keys()))
        raise KeyError(
            f"Unknown placeholder '{placeholder_name}' for layout '{layout_name}'. "
            f"Available: {available}"
        )
    
    return layout_placeholders[placeholder_name]
