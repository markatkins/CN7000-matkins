"""
Progress report slide helpers for Cornelis PowerPoint generation.

Usage:
    from utilities.pptx_helper import create_presentation, save_presentation
    from utilities.pptx_helper.progress_report import (
        add_title_slide, add_status_summary_slide, add_item_list_slide,
        add_comparison_slide, add_section_header
    )
    
    prs = create_presentation()
    add_title_slide(prs, "Q4 Report", "Engineering Status", "John Doe", "Engineer")
    add_status_summary_slide(prs, "Sprint Status", open_count=5, closed_count=12)
    add_item_list_slide(prs, "Open Items", ["Bug #123", "Feature #456"])
    save_presentation(prs, "report.pptx")

Design:
    Each function adds a single slide to the presentation using the appropriate
    layout from the Cornelis template. Functions return the created slide object
    for further customization if needed.
"""

from typing import List, Optional, TYPE_CHECKING

from utilities.pptx_helper.layouts import (
    COVER, CONTENT, CONTENT_WITH_SUBTITLE, TWO_COLUMN, SECTION_HEADER,
    get_placeholder_idx
)

if TYPE_CHECKING:
    from pptx import Presentation
    from pptx.slide import Slide


def add_title_slide(
    prs: "Presentation",
    title: str,
    subtitle: str = "",
    presenter_name: str = "",
    presenter_info: str = ""
) -> "Slide":
    """
    Add a title/cover slide to the presentation.
    
    Uses Layout 1 (Cover with motion) which includes placeholders for
    title, subtitle, and presenter information.
    
    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle text (optional)
        presenter_name: Presenter's name (optional)
        presenter_info: Presenter's title/role (optional)
        
    Returns:
        The created Slide object
        
    Example:
        >>> prs = create_presentation()
        >>> add_title_slide(prs, "Q4 Progress Report", "Engineering Status",
        ...                 "Jane Smith", "Senior Engineer")
    """
    layout = prs.slide_layouts[COVER]
    slide = prs.slides.add_slide(layout)
    
    # Title (idx=0)
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Subtitle (idx=1)
    if subtitle and 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = subtitle
    
    # Presenter name (idx=11)
    if presenter_name and 11 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[11].text = presenter_name
    
    # Presenter info (idx=12)
    if presenter_info and 12 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[12].text = presenter_info
    
    return slide


def add_status_summary_slide(
    prs: "Presentation",
    title: str,
    open_count: int = 0,
    closed_count: int = 0,
    notes: str = "",
    subtitle: str = ""
) -> "Slide":
    """
    Add a status summary slide showing open/closed item counts.
    
    Uses Layout 9 (Title subtitle and content) for a clean summary view.
    
    Args:
        prs: Presentation object
        title: Slide title (e.g., "Sprint Status")
        open_count: Number of open items
        closed_count: Number of closed items
        notes: Additional notes or context
        subtitle: Optional subtitle for context
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_status_summary_slide(prs, "Sprint 42 Status",
        ...                          open_count=5, closed_count=12,
        ...                          notes="On track for release")
    """
    layout = prs.slide_layouts[CONTENT_WITH_SUBTITLE]
    slide = prs.slides.add_slide(layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Subtitle (idx=13)
    if subtitle and 13 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[13].text = subtitle
    
    # Body content (idx=1)
    total = open_count + closed_count
    completion = (closed_count / total * 100) if total > 0 else 0
    
    content_lines = [
        f"Open Items: {open_count}",
        f"Closed Items: {closed_count}",
        f"Total: {total}",
        f"Completion: {completion:.0f}%",
    ]
    
    if notes:
        content_lines.append("")
        content_lines.append(f"Notes: {notes}")
    
    if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = "\n".join(content_lines)
    
    return slide


def add_item_list_slide(
    prs: "Presentation",
    title: str,
    items: List[str],
    item_type: str = "open",
    subtitle: str = ""
) -> "Slide":
    """
    Add a slide with a bullet list of items.
    
    Uses Layout 8 (title and 1 content) for standard bullet lists.
    
    Args:
        prs: Presentation object
        title: Slide title (e.g., "Open Items", "Completed Tasks")
        items: List of item descriptions
        item_type: Type indicator ("open" or "closed") - affects formatting
        subtitle: Optional subtitle
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_item_list_slide(prs, "Open Items",
        ...                     ["Bug #123: Login fails", "Feature #456: Add export"],
        ...                     item_type="open")
    """
    layout = prs.slide_layouts[CONTENT]
    slide = prs.slides.add_slide(layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Format items as bullet points
    prefix = "○" if item_type == "open" else "●"
    formatted_items = [f"{prefix} {item}" for item in items]
    
    # Body content (idx=1)
    if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = "\n".join(formatted_items)
    
    return slide


def add_comparison_slide(
    prs: "Presentation",
    title: str,
    left_title: str,
    left_items: List[str],
    right_title: str,
    right_items: List[str]
) -> "Slide":
    """
    Add a two-column comparison slide (e.g., Plan vs Actual).
    
    Uses Layout 13 (Title and Content - 2 Column) for side-by-side comparison.
    
    Args:
        prs: Presentation object
        title: Slide title (e.g., "Plan vs Actual")
        left_title: Header for left column
        left_items: List of items for left column
        right_title: Header for right column
        right_items: List of items for right column
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_comparison_slide(prs, "Sprint Goals vs Delivered",
        ...                      "Planned", ["Feature A", "Feature B"],
        ...                      "Delivered", ["Feature A ✓", "Feature B (80%)"])
    """
    layout = prs.slide_layouts[TWO_COLUMN]
    slide = prs.slides.add_slide(layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Left column (idx=1)
    left_content = f"{left_title}:\n" + "\n".join(f"• {item}" for item in left_items)
    if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = left_content
    
    # Right column (idx=11)
    right_content = f"{right_title}:\n" + "\n".join(f"• {item}" for item in right_items)
    if 11 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[11].text = right_content
    
    return slide


def add_section_header(
    prs: "Presentation",
    title: str,
    subtitle: str = ""
) -> "Slide":
    """
    Add a section header/divider slide.
    
    Uses Layout 37 (Section Header) for visual separation between sections.
    
    Args:
        prs: Presentation object
        title: Section title
        subtitle: Optional subtitle or description
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_section_header(prs, "Technical Details", "Architecture and Implementation")
    """
    layout = prs.slide_layouts[SECTION_HEADER]
    slide = prs.slides.add_slide(layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Subtitle (idx=1)
    if subtitle and 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = subtitle
    
    return slide
