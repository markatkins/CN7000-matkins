"""
Technical presentation slide helpers for Cornelis PowerPoint generation.

Usage:
    from utilities.pptx_helper import create_presentation, save_presentation
    from utilities.pptx_helper.technical import (
        add_content_slide, add_image_slide, add_image_content_slide,
        add_table_slide, add_two_column_slide, add_code_slide
    )
    
    prs = create_presentation()
    add_content_slide(prs, "Architecture", ["Component A", "Component B"])
    add_table_slide(prs, "Metrics", ["Name", "Value"], [["Latency", "10ms"]])
    save_presentation(prs, "tech.pptx")

Design:
    Each function adds a single slide to the presentation using the appropriate
    layout from the Cornelis template. Functions return the created slide object
    for further customization if needed.
"""

from pathlib import Path
from typing import List, Optional, Union, TYPE_CHECKING

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import MSO_ANCHOR

from utilities.pptx_helper.layouts import (
    CONTENT, CONTENT_WITH_SUBTITLE, TWO_COLUMN, CONTENT_AND_PICTURE,
    TABLE, BLANK
)

if TYPE_CHECKING:
    from pptx import Presentation
    from pptx.slide import Slide

MAX_LINES_PER_SLIDE = 19
TABLE_FONT_SIZE_DEFAULT = Pt(18)
TABLE_FONT_SIZE_MIN = Pt(14)
TITLE_FONT_SIZE_DEFAULT = Pt(36)
TITLE_FONT_SIZE_MIN = Pt(24)
TITLE_PAREN_FONT_SIZE = Pt(20)
CELL_MARGIN_NARROW = Emu(45720)
TABLE_TOP_POSITION = Inches(1.1)
TABLE_WIDTH = Inches(12.3)
CHARS_PER_INCH_14PT = 12


def add_content_slide(
    prs: "Presentation",
    title: str,
    bullets: List[str],
    subtitle: str = ""
) -> "Slide":
    """
    Add a standard content slide with bullet points.
    
    Uses Layout 8 (title and 1 content) or Layout 9 if subtitle provided.
    
    Args:
        prs: Presentation object
        title: Slide title
        bullets: List of bullet point strings
        subtitle: Optional subtitle (uses different layout if provided)
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_content_slide(prs, "Key Features",
        ...                   ["High performance", "Easy integration", "Scalable"])
    """
    layout_idx = CONTENT_WITH_SUBTITLE if subtitle else CONTENT
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Subtitle (idx=13) - only for CONTENT_WITH_SUBTITLE layout
    if subtitle and 13 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[13].text = subtitle
    
    # Body content (idx=1)
    content = "\n".join(f"• {bullet}" for bullet in bullets)
    if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = content
    
    return slide


def add_image_slide(
    prs: "Presentation",
    title: str,
    image_path: Union[str, Path],
    caption: str = ""
) -> "Slide":
    """
    Add a slide with a prominent image.
    
    Uses Layout 25 (1/2 Content and Picture Insert) with image taking focus.
    
    Args:
        prs: Presentation object
        title: Slide title
        image_path: Path to image file (PNG, JPG, etc.)
        caption: Optional caption text
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_image_slide(prs, "System Architecture", "diagrams/arch.png",
        ...                 caption="High-level component diagram")
    """
    layout = prs.slide_layouts[CONTENT_AND_PICTURE]
    slide = prs.slides.add_slide(layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Caption in body (idx=1)
    if caption and 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = caption
    
    # Image (idx=13)
    if 13 in [p.placeholder_format.idx for p in slide.placeholders]:
        picture_placeholder = slide.placeholders[13]
        picture_placeholder.insert_picture(str(image_path))
    
    return slide


def add_image_content_slide(
    prs: "Presentation",
    title: str,
    content: str,
    image_path: Union[str, Path]
) -> "Slide":
    """
    Add a slide with text content and an accompanying image.
    
    Uses Layout 25 (1/2 Content and Picture Insert) for balanced layout.
    
    Args:
        prs: Presentation object
        title: Slide title
        content: Text content (can include newlines for multiple paragraphs)
        image_path: Path to image file
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_image_content_slide(prs, "Data Flow",
        ...                         "Key components:\\n• API Gateway\\n• Service Mesh",
        ...                         "diagrams/dataflow.png")
    """
    layout = prs.slide_layouts[CONTENT_AND_PICTURE]
    slide = prs.slides.add_slide(layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Content (idx=1)
    if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = content
    
    # Image (idx=13)
    if 13 in [p.placeholder_format.idx for p in slide.placeholders]:
        picture_placeholder = slide.placeholders[13]
        picture_placeholder.insert_picture(str(image_path))
    
    return slide


def _set_title_with_auto_size(shape, title: str) -> None:
    """Set title text with automatic font sizing for long titles."""
    if not shape:
        return
    
    shape.text = title
    
    title_len = len(title)
    has_paren = '(' in title and ')' in title
    
    if title_len > 60:
        font_size = TITLE_FONT_SIZE_MIN
    elif title_len > 45:
        font_size = Pt(28)
    elif title_len > 35:
        font_size = Pt(32)
    else:
        font_size = TITLE_FONT_SIZE_DEFAULT
    
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
        
        if has_paren and title_len > 30:
            paren_start = title.find('(')
            if paren_start > 0:
                shape.text = ""
                p = shape.text_frame.paragraphs[0]
                main_run = p.add_run()
                main_run.text = title[:paren_start]
                main_run.font.size = font_size
                
                paren_run = p.add_run()
                paren_run.text = title[paren_start:]
                paren_run.font.size = TITLE_PAREN_FONT_SIZE
    except Exception:
        pass


def _apply_table_formatting(table, num_rows: int, num_cols: int, font_size) -> None:
    for row_idx in range(num_rows):
        for col_idx in range(num_cols):
            cell = table.cell(row_idx, col_idx)
            cell.margin_left = CELL_MARGIN_NARROW
            cell.margin_right = CELL_MARGIN_NARROW
            cell.margin_top = CELL_MARGIN_NARROW
            cell.margin_bottom = CELL_MARGIN_NARROW
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            try:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = font_size
                        if row_idx == 0:
                            run.font.bold = True
            except Exception:
                pass


def _calculate_column_widths(headers: List[str], rows: List[List[str]], num_cols: int) -> List[int]:
    max_chars = [len(h) for h in headers]
    
    for row in rows:
        for col_idx, cell_value in enumerate(row):
            if col_idx < num_cols:
                cell_str = str(cell_value)
                lines = cell_str.split('\n')
                max_line_len = max(len(line) for line in lines) if lines else 0
                max_chars[col_idx] = max(max_chars[col_idx], max_line_len)
    
    total_chars = sum(max_chars)
    if total_chars == 0:
        return [Inches(TABLE_WIDTH.inches / num_cols)] * num_cols
    
    col_widths = []
    for chars in max_chars:
        proportion = chars / total_chars
        width_inches = max(0.8, proportion * TABLE_WIDTH.inches)
        col_widths.append(Inches(width_inches))
    
    total_width = sum(w.inches for w in col_widths)
    if total_width > TABLE_WIDTH.inches:
        scale = TABLE_WIDTH.inches / total_width
        col_widths = [Inches(w.inches * scale) for w in col_widths]
    
    return col_widths


def _count_lines_in_row(row: List[str], col_widths: List[int], num_cols: int) -> int:
    max_lines = 1
    for col_idx, cell_value in enumerate(row):
        if col_idx >= num_cols:
            break
        cell_str = str(cell_value)
        explicit_lines = cell_str.count('\n') + 1
        
        col_width_inches = col_widths[col_idx].inches if col_idx < len(col_widths) else 1.0
        chars_per_line = int(col_width_inches * CHARS_PER_INCH_14PT)
        chars_per_line = max(chars_per_line, 5)
        
        for line in cell_str.split('\n'):
            wrapped_lines = (len(line) + chars_per_line - 1) // chars_per_line if line else 1
            explicit_lines += wrapped_lines - 1
        
        max_lines = max(max_lines, explicit_lines)
    
    return max_lines


def _split_rows_by_lines(headers: List[str], rows: List[List[str]], col_widths: List[int], num_cols: int) -> List[List[List[str]]]:
    chunks = []
    current_chunk = []
    current_lines = 1
    
    for row in rows:
        row_lines = _count_lines_in_row(row, col_widths, num_cols)
        
        if current_lines + row_lines > MAX_LINES_PER_SLIDE and current_chunk:
            chunks.append(current_chunk)
            current_chunk = []
            current_lines = 1
        
        current_chunk.append(row)
        current_lines += row_lines
    
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks


def add_table_slide(
    prs: "Presentation",
    title: str,
    headers: List[str],
    rows: List[List[str]],
    description: str = ""
) -> "Slide":
    slides_created = []
    total_rows = len(rows)
    num_cols = len(headers)
    
    col_widths = _calculate_column_widths(headers, rows, num_cols)
    chunks = _split_rows_by_lines(headers, rows, col_widths, num_cols)
    total_pages = len(chunks)
    
    if total_rows > 15:
        font_size = TABLE_FONT_SIZE_MIN
    elif total_rows > 10:
        font_size = Pt(16)
    else:
        font_size = TABLE_FONT_SIZE_DEFAULT
    
    for page_num, chunk_rows in enumerate(chunks, start=1):
        if total_pages > 1:
            slide_title = f"{title} ({page_num}/{total_pages})"
        else:
            slide_title = title
        
        slide = _create_single_table_slide(prs, slide_title, headers, chunk_rows, col_widths, font_size)
        slides_created.append(slide)
    
    return slides_created[0] if slides_created else None


def _create_single_table_slide(
    prs: "Presentation",
    title: str,
    headers: List[str],
    rows: List[List[str]],
    col_widths: List[int],
    font_size
) -> "Slide":
    layout = prs.slide_layouts[CONTENT]
    slide = prs.slides.add_slide(layout)
    
    _set_title_with_auto_size(slide.shapes.title, title)
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    left = Inches(0.5)
    top = TABLE_TOP_POSITION
    width = TABLE_WIDTH
    
    row_height_inches = 0.28 if font_size <= TABLE_FONT_SIZE_MIN else 0.32
    height = Inches(min(num_rows * row_height_inches, 5.8))
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    for col_idx in range(num_cols):
        if col_idx < len(col_widths):
            table.columns[col_idx].width = col_widths[col_idx]
    
    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header
    
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_value in enumerate(row_data):
            if col_idx < num_cols:
                table.cell(row_idx, col_idx).text = str(cell_value)
    
    _apply_table_formatting(table, num_rows, num_cols, font_size)
    
    return slide


def add_two_column_slide(
    prs: "Presentation",
    title: str,
    left_content: str,
    right_content: str
) -> "Slide":
    """
    Add a two-column slide for comparisons or parallel content.
    
    Uses Layout 13 (Title and Content - 2 Column).
    
    Args:
        prs: Presentation object
        title: Slide title
        left_content: Content for left column (can include newlines)
        right_content: Content for right column (can include newlines)
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_two_column_slide(prs, "Before vs After",
        ...                      left_content="Old approach:\\n- Slow\\n- Complex",
        ...                      right_content="New approach:\\n- Fast\\n- Simple")
    """
    layout = prs.slide_layouts[TWO_COLUMN]
    slide = prs.slides.add_slide(layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Left column (idx=1)
    if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[1].text = left_content
    
    # Right column (idx=11)
    if 11 in [p.placeholder_format.idx for p in slide.placeholders]:
        slide.placeholders[11].text = right_content
    
    return slide


def add_code_slide(
    prs: "Presentation",
    title: str,
    code: str,
    language: str = ""
) -> "Slide":
    """
    Add a slide with a code block.
    
    Uses Layout 8 (title and 1 content) with monospace formatting.
    Note: python-pptx doesn't support syntax highlighting, so code
    is displayed as plain monospace text.
    
    Args:
        prs: Presentation object
        title: Slide title
        code: Code string (preserves whitespace and newlines)
        language: Optional language name for display (informational only)
        
    Returns:
        The created Slide object
        
    Example:
        >>> add_code_slide(prs, "Example Usage",
        ...                code="def hello():\\n    print('Hello, World!')",
        ...                language="Python")
    """
    layout = prs.slide_layouts[CONTENT]
    slide = prs.slides.add_slide(layout)
    
    # Title (include language if provided)
    title_text = f"{title} ({language})" if language else title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    
    # Code in body (idx=1)
    if 1 in [p.placeholder_format.idx for p in slide.placeholders]:
        placeholder = slide.placeholders[1]
        placeholder.text = code
        
        # Try to set monospace font on all runs
        try:
            for paragraph in placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Consolas"
                    run.font.size = Pt(12)
        except Exception:
            # If font setting fails, continue with default
            pass
    
    return slide
