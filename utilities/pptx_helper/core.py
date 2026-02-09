"""
Core functionality for Cornelis PowerPoint generation.

Usage:
    from utilities.pptx_helper import load_template, create_presentation, save_presentation
    
    # Create a new presentation from template
    prs = create_presentation()
    
    # Add slides...
    
    # Save
    save_presentation(prs, "output/my_presentation.pptx")

Design:
    Provides template loading with .potx workaround, presentation creation,
    and saving utilities. The .potx workaround modifies the content type
    in the template's [Content_Types].xml to make python-pptx accept it.

Traceability:
    Template: templates/Standard PPT Template_Light.potx
    Spec: templates/template-spec.md
"""

import logging
import os
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Optional, Union

from pptx import Presentation

logger = logging.getLogger(__name__)

# Default template path (relative to project root)
DEFAULT_TEMPLATE = "templates/Standard PPT Template_Light.potx"

# Expected number of layouts in the template (for validation)
EXPECTED_LAYOUT_COUNT = 39


def _find_template_path(template_path: Optional[str] = None) -> Path:
    """
    Find the template file path.
    
    Searches in order:
    1. Provided path (if given)
    2. Relative to current working directory
    3. Relative to this module's location
    
    Args:
        template_path: Optional explicit path to template
        
    Returns:
        Path to template file
        
    Raises:
        FileNotFoundError: If template cannot be found
    """
    if template_path:
        path = Path(template_path)
        if path.exists():
            return path
        raise FileNotFoundError(f"Template not found: {template_path}")
    
    # Try relative to cwd
    cwd_path = Path.cwd() / DEFAULT_TEMPLATE
    if cwd_path.exists():
        return cwd_path
    
    # Try relative to module
    module_dir = Path(__file__).parent.parent.parent
    module_path = module_dir / DEFAULT_TEMPLATE
    if module_path.exists():
        return module_path
    
    raise FileNotFoundError(
        f"Template not found. Searched:\n"
        f"  - {cwd_path}\n"
        f"  - {module_path}\n"
        f"Provide explicit path via template_path parameter."
    )


def _convert_potx_to_pptx(potx_path: Path) -> str:
    """
    Convert a .potx file to .pptx by modifying the content type.
    
    The .potx format uses a different content type in [Content_Types].xml.
    This function creates a temporary copy with the content type changed
    to make python-pptx accept it.
    
    Args:
        potx_path: Path to the .potx file
        
    Returns:
        Path to temporary .pptx file (caller must clean up)
    """
    # Create temp directory for extraction
    temp_dir = tempfile.mkdtemp()
    temp_pptx = tempfile.mktemp(suffix='.pptx')
    
    try:
        # Extract the potx
        with zipfile.ZipFile(potx_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Modify [Content_Types].xml
        content_types_path = os.path.join(temp_dir, '[Content_Types].xml')
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Replace template content type with presentation content type
        content = content.replace(
            'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
        )
        
        with open(content_types_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        # Repack as pptx
        with zipfile.ZipFile(temp_pptx, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zipf.write(file_path, arcname)
        
        return temp_pptx
        
    finally:
        # Clean up temp directory
        shutil.rmtree(temp_dir, ignore_errors=True)


def load_template(template_path: Optional[str] = None) -> Presentation:
    """
    Load the PowerPoint template.
    
    Handles .potx files by converting content type to .pptx format,
    as python-pptx rejects .potx content type.
    
    Args:
        template_path: Optional path to template file. If not provided,
                      searches for default template.
                      
    Returns:
        Presentation object loaded from template
        
    Raises:
        FileNotFoundError: If template file not found
        ValueError: If template doesn't have expected number of layouts
        
    Example:
        >>> prs = load_template()
        >>> print(f"Loaded {len(prs.slide_layouts)} layouts")
    """
    path = _find_template_path(template_path)
    logger.info(f"Loading template from: {path}")
    
    temp_pptx = None
    try:
        # Check if it's a .potx file that needs conversion
        if path.suffix.lower() == '.potx':
            logger.debug("Converting .potx to .pptx format")
            temp_pptx = _convert_potx_to_pptx(path)
            prs = Presentation(temp_pptx)
        else:
            prs = Presentation(str(path))
        
        # Validate template
        layout_count = len(prs.slide_layouts)
        if layout_count != EXPECTED_LAYOUT_COUNT:
            logger.warning(
                f"Template has {layout_count} layouts, expected {EXPECTED_LAYOUT_COUNT}. "
                f"Layout indices may not match documentation."
            )
        else:
            logger.debug(f"Template validated: {layout_count} layouts")
        
        return prs
        
    finally:
        # Clean up temp file
        if temp_pptx and os.path.exists(temp_pptx):
            os.remove(temp_pptx)


def create_presentation(template_path: Optional[str] = None) -> Presentation:
    """
    Create a new presentation from the Cornelis template.
    
    This is the main entry point for creating presentations.
    
    Args:
        template_path: Optional path to template file
        
    Returns:
        New Presentation object ready for adding slides
        
    Example:
        >>> prs = create_presentation()
        >>> slide = prs.slides.add_slide(prs.slide_layouts[8])
        >>> slide.shapes.title.text = "My Title"
    """
    return load_template(template_path)


def save_presentation(
    prs: Presentation,
    output_path: Union[str, Path],
    create_dirs: bool = True
) -> Path:
    """
    Save a presentation to file.
    
    Args:
        prs: Presentation object to save
        output_path: Path for output file (should end in .pptx)
        create_dirs: If True, create parent directories if needed
        
    Returns:
        Path to saved file
        
    Raises:
        ValueError: If output path doesn't end in .pptx
        
    Example:
        >>> prs = create_presentation()
        >>> # ... add slides ...
        >>> path = save_presentation(prs, "output/report.pptx")
        >>> print(f"Saved to: {path}")
    """
    output_path = Path(output_path)
    
    if output_path.suffix.lower() != '.pptx':
        raise ValueError(f"Output path must end in .pptx, got: {output_path}")
    
    if create_dirs:
        output_path.parent.mkdir(parents=True, exist_ok=True)
    
    logger.info(f"Saving presentation to: {output_path}")
    prs.save(str(output_path))
    
    return output_path
