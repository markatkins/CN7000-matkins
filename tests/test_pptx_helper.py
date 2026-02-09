"""
Tests for the pptx_helper module.

Run with: pytest tests/test_pptx_helper.py -v
"""

import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor


class TestTemplateLoading:
    """Tests for template loading functionality."""
    
    def test_template_loads(self):
        """Verify template loads without errors."""
        from utilities.pptx_helper import load_template
        
        prs = load_template()
        assert prs is not None
    
    def test_template_has_expected_layouts(self):
        """Verify template has 39 layouts."""
        from utilities.pptx_helper import load_template
        
        prs = load_template()
        assert len(prs.slide_layouts) == 39
    
    def test_create_presentation(self):
        """Verify create_presentation returns valid presentation."""
        from utilities.pptx_helper import create_presentation
        
        prs = create_presentation()
        assert prs is not None
        assert len(prs.slide_layouts) == 39


class TestLayouts:
    """Tests for layout constants and helpers."""
    
    def test_layouts_defined(self):
        """Verify layout constants are defined."""
        from utilities.pptx_helper.layouts import (
            LAYOUTS, CONTENT, TWO_COLUMN, TABLE, SECTION_HEADER
        )
        
        assert 'content' in LAYOUTS
        assert CONTENT == 8
        assert TWO_COLUMN == 13
        assert TABLE == 30
        assert SECTION_HEADER == 37
    
    def test_get_layout(self):
        """Verify get_layout returns correct layout."""
        from utilities.pptx_helper import create_presentation, get_layout
        
        prs = create_presentation()
        layout = get_layout(prs, "content")
        assert layout is not None
    
    def test_get_layout_invalid_name(self):
        """Verify get_layout raises KeyError for invalid name."""
        from utilities.pptx_helper import create_presentation, get_layout
        
        prs = create_presentation()
        with pytest.raises(KeyError):
            get_layout(prs, "nonexistent_layout")


class TestColors:
    """Tests for color constants."""
    
    def test_colors_defined(self):
        """Verify color constants are defined correctly."""
        from utilities.pptx_helper.colors import (
            PRIMARY_PURPLE, CORNELIS_COLORS
        )
        
        assert PRIMARY_PURPLE == RGBColor(100, 0, 185)
        assert 'primary_purple' in CORNELIS_COLORS
        assert CORNELIS_COLORS['primary_purple'] == PRIMARY_PURPLE
    
    def test_all_colors_are_rgbcolor(self):
        """Verify all colors in CORNELIS_COLORS are RGBColor objects."""
        from utilities.pptx_helper.colors import CORNELIS_COLORS
        
        for name, color in CORNELIS_COLORS.items():
            assert isinstance(color, RGBColor), f"{name} is not RGBColor"


class TestProgressReportGeneration:
    """Tests for progress report slide generation."""
    
    def test_progress_report_imports(self):
        """Verify progress report functions import successfully."""
        from utilities.pptx_helper.progress_report import (
            add_title_slide, add_status_summary_slide, add_item_list_slide,
            add_comparison_slide, add_section_header
        )
        
        assert callable(add_title_slide)
        assert callable(add_status_summary_slide)
        assert callable(add_item_list_slide)
        assert callable(add_comparison_slide)
        assert callable(add_section_header)
    
    def test_progress_report_generation(self):
        """Test end-to-end progress report generation."""
        from utilities.pptx_helper import create_presentation, save_presentation
        from utilities.pptx_helper.progress_report import (
            add_title_slide, add_status_summary_slide, add_item_list_slide
        )
        
        prs = create_presentation()
        initial_slides = len(prs.slides)
        
        add_title_slide(prs, 'Q4 Progress Report', 'Engineering Status', 'Test User', 'Engineer')
        add_status_summary_slide(prs, 'Sprint Status', open_count=5, closed_count=12, notes='On track')
        add_item_list_slide(prs, 'Open Items', ['Item A', 'Item B', 'Item C'], item_type='open')
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name
        
        try:
            save_presentation(prs, output_path)
            
            # Verify by re-opening
            prs2 = Presentation(output_path)
            added_slides = len(prs2.slides) - initial_slides
            assert added_slides == 3
        finally:
            os.unlink(output_path)


class TestTechnicalPresentationGeneration:
    """Tests for technical presentation slide generation."""
    
    def test_technical_imports(self):
        """Verify technical presentation functions import successfully."""
        from utilities.pptx_helper.technical import (
            add_content_slide, add_image_slide, add_image_content_slide,
            add_table_slide, add_two_column_slide, add_code_slide
        )
        
        assert callable(add_content_slide)
        assert callable(add_image_slide)
        assert callable(add_image_content_slide)
        assert callable(add_table_slide)
        assert callable(add_two_column_slide)
        assert callable(add_code_slide)
    
    def test_technical_presentation_generation(self):
        """Test end-to-end technical presentation generation."""
        from utilities.pptx_helper import create_presentation, save_presentation
        from utilities.pptx_helper.technical import (
            add_content_slide, add_table_slide, add_two_column_slide
        )
        
        prs = create_presentation()
        initial_slides = len(prs.slides)
        
        add_content_slide(prs, 'Architecture Overview', ['Component A', 'Component B', 'Component C'])
        add_table_slide(prs, 'Performance Metrics', 
            headers=['Metric', 'Target', 'Actual'],
            rows=[['Latency', '10ms', '8ms'], ['Throughput', '1000/s', '1200/s']])
        add_two_column_slide(prs, 'Before vs After',
            left_content='Old approach:\n- Slow\n- Complex',
            right_content='New approach:\n- Fast\n- Simple')
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name
        
        try:
            save_presentation(prs, output_path)
            
            # Verify by re-opening
            prs2 = Presentation(output_path)
            added_slides = len(prs2.slides) - initial_slides
            assert added_slides == 3
        finally:
            os.unlink(output_path)


class TestGeneratedPptxValid:
    """Tests that generated PPTX files are valid."""
    
    def test_generated_pptx_can_be_reopened(self):
        """Verify generated PPTX files can be re-opened without errors."""
        from utilities.pptx_helper import create_presentation, save_presentation
        from utilities.pptx_helper.progress_report import add_title_slide
        
        prs = create_presentation()
        add_title_slide(prs, 'Test', 'Subtitle', 'Name', 'Info')
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name
        
        try:
            save_presentation(prs, output_path)
            
            # Re-open should not raise
            prs2 = Presentation(output_path)
            assert prs2 is not None
        finally:
            os.unlink(output_path)
    
    def test_save_creates_parent_directories(self):
        """Verify save_presentation creates parent directories."""
        from utilities.pptx_helper import create_presentation, save_presentation
        
        prs = create_presentation()
        
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / 'subdir' / 'nested' / 'test.pptx'
            save_presentation(prs, output_path)
            assert output_path.exists()


class TestModuleDocumentation:
    """Tests for module documentation."""
    
    def test_module_has_docstring(self):
        """Verify module has docstring with examples."""
        from utilities import pptx_helper
        
        assert pptx_helper.__doc__ is not None
        # Check for key content
        assert 'Quick Start' in pptx_helper.__doc__ or 'Example' in pptx_helper.__doc__


if __name__ == '__main__':
    pytest.main([__file__, '-v'])
